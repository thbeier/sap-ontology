"""Generate docs/sap-ontology-overview.pptx — the client-facing overview deck.

Usage (from repo root, with .venv active):
    python docs/generate_overview_pptx.py

Re-run whenever README / CHANGELOG / class catalog drift; the deck is built
from constants in this script, not from parsing the docs.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

DECK_PATH = Path(__file__).resolve().parent / "sap-ontology-overview.pptx"
VERSION = "v0.2.0"

NAVY = RGBColor(0x0B, 0x1F, 0x3A)
ACCENT = RGBColor(0x00, 0x6E, 0xC7)
GREY = RGBColor(0x4A, 0x4A, 0x4A)
LIGHT = RGBColor(0xF2, 0xF5, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _set_text(frame, text, *, size=18, bold=False, color=GREY, align=PP_ALIGN.LEFT):
    frame.clear()
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_bullets(frame, bullets, *, size=16):
    frame.clear()
    for i, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = bullet
        run.font.size = Pt(size)
        run.font.color.rgb = GREY
        p.level = 0


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def _add_title_band(slide, prs, title, subtitle=None):
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1)
    )
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    _set_text(band.text_frame, title, size=28, bold=True, color=WHITE)
    band.text_frame.margin_left = Inches(0.4)
    band.text_frame.margin_top = Inches(0.25)
    if subtitle:
        sub = slide.shapes.add_textbox(
            Inches(0.4), Inches(1.15), prs.slide_width - Inches(0.8), Inches(0.4)
        )
        _set_text(sub.text_frame, subtitle, size=14, color=ACCENT, bold=True)


def _add_footer(slide, prs, page, total):
    fb = slide.shapes.add_textbox(
        Inches(0.4),
        prs.slide_height - Inches(0.4),
        prs.slide_width - Inches(0.8),
        Inches(0.3),
    )
    _set_text(
        fb.text_frame,
        f"SAP Ontology for Context Engineering · {VERSION} · {page}/{total}",
        size=10,
        color=GREY,
    )


# ---------- slide builders ----------

def slide_title(prs):
    s = _blank_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY

    title = s.shapes.add_textbox(
        Inches(0.8), Inches(2.2), prs.slide_width - Inches(1.6), Inches(1.6)
    )
    _set_text(
        title.text_frame,
        "SAP Ontology for\nContext Engineering",
        size=44,
        bold=True,
        color=WHITE,
    )

    sub = s.shapes.add_textbox(
        Inches(0.8), Inches(4.0), prs.slide_width - Inches(1.6), Inches(0.6)
    )
    _set_text(
        sub.text_frame,
        "A technology-independent semantic model for reasoning across "
        "Architecture × Business Process × Implementation / Change.",
        size=18,
        color=ACCENT,
    )

    foot = s.shapes.add_textbox(
        Inches(0.8), prs.slide_height - Inches(0.9), prs.slide_width - Inches(1.6), Inches(0.5)
    )
    _set_text(foot.text_frame, f"{VERSION} · CC-BY-SA 4.0", size=12, color=WHITE)
    return s


def slide_problem(prs):
    s = _blank_slide(prs)
    _add_title_band(s, prs, "The problem")
    body = s.shapes.add_textbox(
        Inches(0.5), Inches(1.6), prs.slide_width - Inches(1.0), Inches(5.0)
    )
    _add_bullets(
        body.text_frame,
        [
            "SAP landscapes carry decades of tribal knowledge — locked in spreadsheets and consultants' heads.",
            "Traditional RAG retrieves text chunks. AI agents need typed traversal: change → config → process → org.",
            "Every project re-discovers the same relations. Reuse across clients is near-zero today.",
            "No common vocabulary means no shared graph, no audit trail, no scenario branching.",
        ],
        size=18,
    )
    return s


def slide_what_it_is(prs):
    s = _blank_slide(prs)
    _add_title_band(s, prs, "What it is")
    body = s.shapes.add_textbox(
        Inches(0.5), Inches(1.6), prs.slide_width - Inches(1.0), Inches(5.0)
    )
    _add_bullets(
        body.text_frame,
        [
            "A Deloitte-curated, technology-independent upper ontology — published under CC-BY-SA 4.0.",
            "JSON-LD 1.1 schema + SHACL shapes + canonical examples — no vendor lock-in.",
            "Anchored on industry standards (ArchiMate, BPMN, APQC, ITIL 4, TOGAF) — not invented from scratch.",
            "Companion proprietary federation runtime: Excel/CSV → SHACL → Neo4j, with per-tenant scenarios.",
            f"Current release: {VERSION} — 20 concrete classes, 30+ inter-domain relations.",
        ],
        size=18,
    )
    return s


def slide_five_domains(prs):
    s = _blank_slide(prs)
    _add_title_band(
        s, prs, "Five domains", subtitle="Three primary + two cross-cutting"
    )

    box_w = Inches(2.85)
    box_h = Inches(2.0)
    top_y = Inches(1.8)
    bot_y = Inches(4.0)
    primary = [
        ("Architecture", "ApplicationComponent · Integration · TBB · DataObject", Inches(0.5), top_y),
        ("Business Process", "Process · Activity · Event · Decision · BusinessDocument", Inches(3.55), top_y),
        ("Implementation / Change", "Configuration · Change · Transport · TestCase · Incident · Requirement", Inches(6.6), top_y),
    ]
    crosscut = [
        ("Organization (cross-cutter)", "OrgUnit · Role · Capability · User (GDPR-gated)", Inches(0.5), bot_y),
        ("Scenario (cross-cutter)", "Scenario with type ∈ {as-is, to-be, variant} and lifecycleState", Inches(6.6), bot_y),
    ]
    for label, body, x, y in primary:
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
        box.line.color.rgb = NAVY
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT
        tf = box.text_frame
        tf.word_wrap = True
        _set_text(tf, label, size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "\n" + body
        run.font.size = Pt(12)
        run.font.color.rgb = GREY

    cross_w = Inches(5.9)
    for label, body, x, y in crosscut:
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cross_w, box_h)
        box.line.color.rgb = ACCENT
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        tf = box.text_frame
        tf.word_wrap = True
        _set_text(tf, label, size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "\n" + body
        run.font.size = Pt(12)
        run.font.color.rgb = GREY

    note = s.shapes.add_textbox(
        Inches(0.5), Inches(6.3), prs.slide_width - Inches(1.0), Inches(0.6)
    )
    _set_text(
        note.text_frame,
        "Every domain instance carries a Provenance (where it came from) and an inScenario "
        "edge (which state it belongs to). Both are non-negotiable, enforced by SHACL.",
        size=12,
        color=GREY,
    )
    return s


def slide_class_catalog(prs):
    s = _blank_slide(prs)
    _add_title_band(s, prs, "Class catalog", subtitle="20 concrete classes + 2 foundation classes")

    rows = [
        ("Architecture", "ApplicationComponent, Integration, TechnologyBuildingBlock, DataObject"),
        ("Business Process", "Process (recursive), Activity, Event, Decision, BusinessDocument"),
        ("Implementation / Change", "Configuration, Change, Transport, TestCase, Incident, Requirement"),
        ("Organization", "OrgUnit, Role, Capability, User (GDPR-gated)"),
        ("Scenario", "Scenario (lifecycleState ∈ {draft, active, locked, superseded})"),
        ("Foundation", "Provenance (mandatory), DomainInstance (abstract marker)"),
    ]

    table_shape = s.shapes.add_table(
        len(rows) + 1, 2, Inches(0.5), Inches(1.6),
        prs.slide_width - Inches(1.0), Inches(4.5),
    )
    tbl = table_shape.table
    tbl.columns[0].width = Inches(2.3)
    tbl.columns[1].width = prs.slide_width - Inches(3.3)

    headers = ["Domain", "Classes"]
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        _set_text(cell.text_frame, h, size=14, bold=True, color=WHITE)
    for r, (dom, cls) in enumerate(rows, start=1):
        c0 = tbl.cell(r, 0)
        _set_text(c0.text_frame, dom, size=12, bold=True, color=NAVY)
        c1 = tbl.cell(r, 1)
        _set_text(c1.text_frame, cls, size=12, color=GREY)
        if r % 2 == 0:
            c0.fill.solid(); c0.fill.fore_color.rgb = LIGHT
            c1.fill.solid(); c1.fill.fore_color.rgb = LIGHT

    note = s.shapes.add_textbox(
        Inches(0.5), Inches(6.3), prs.slide_width - Inches(1.0), Inches(0.6)
    )
    _set_text(
        note.text_frame,
        f"{VERSION} added sap:transactionCode (Activity) and sap:configurationTransaction "
        "(Configuration) as canonical SAP T-code / SPRO handles.",
        size=12, color=GREY,
    )
    return s


def slide_inter_domain(prs):
    s = _blank_slide(prs)
    _add_title_band(
        s, prs, "Inter-domain relations — the differentiator",
        subtitle="Why this is more than a class diagram",
    )
    body = s.shapes.add_textbox(
        Inches(0.5), Inches(1.6), prs.slide_width - Inches(1.0), Inches(4.9)
    )
    _add_bullets(
        body.text_frame,
        [
            "Activity      ─ supportedBy       →  ApplicationComponent       (which app runs this step)",
            "Activity      ─ realizedBy        →  Configuration              (which customizing implements the rule)",
            "BusinessDocument ─ represents     →  DataObject                 (document ↔ system-of-record entity)",
            "Change        ─ affects           →  Configuration              (what does this change touch)",
            "Change        ─ impacts           →  Activity                   (DERIVED from affects + realizedBy — query-time)",
            "Incident      ─ linkedTo          →  Activity / App / Config    (AMS entry point)",
            "TestCase      ─ validates         →  Change / Config / Activity (test coverage chain)",
            "Role          ─ executes          →  Activity                   (RACI Responsible — the org link)",
            "Role          ─ authorizedFor     →  Configuration              (SAP authorization-model link)",
            "OrgUnit       ─ owns              →  Configuration / Process    (custody / approval chain)",
            "OrgUnit       ─ consumes          →  ApplicationComponent       (which org uses which system)",
            "Process       ─ realizesCapability → Capability                 (TOGAF / BIZBOK capability map)",
            "Requirement   ─ targetsScenario   →  Scenario                   (target-state binding)",
            "Requirement   ─ implementedBy     →  Change                     (delivery chain)",
        ],
        size=12,
    )
    note = s.shapes.add_textbox(
        Inches(0.5), Inches(6.4), prs.slide_width - Inches(1.0), Inches(0.5)
    )
    _set_text(
        note.text_frame,
        "These edges are the queries clients want to ask: \"if I change X, what activities break?\" "
        "An ontology without inter-domain edges is just a glossary.",
        size=12, color=ACCENT, bold=True,
    )
    return s


def slide_artifact_bridge(prs):
    """Documents ↔ data — a concrete traversal making BusinessDocument's role tangible."""
    s = _blank_slide(prs)
    _add_title_band(
        s, prs, "Documents ↔ data — the business artifact bridge",
        subtitle="One worked traversal across all three primary domains",
    )

    # Four entity boxes laid out left-to-right
    box_w = Inches(2.5)
    box_h = Inches(1.4)
    y = Inches(1.9)
    gap = Inches(0.6)
    left_margin = Inches(0.45)

    entities = [
        ("Activity",            "Create Sales Order\nT-code: VA01",         NAVY,   LIGHT),
        ("BusinessDocument",    "Sales Order doc",                          NAVY,   LIGHT),
        ("DataObject",          "Customer · SalesOrder",                    NAVY,   LIGHT),
        ("ApplicationComponent","SAP S/4 — Sales & Distribution",           ACCENT, WHITE),
    ]

    box_xs = []
    for i, (label, body_text, border, fill) in enumerate(entities):
        x = left_margin + (box_w + gap) * i
        box_xs.append(x)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
        box.line.color.rgb = border
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        tf = box.text_frame
        tf.word_wrap = True
        _set_text(tf, label, size=14, bold=True, color=border, align=PP_ALIGN.CENTER)
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "\n" + body_text
        run.font.size = Pt(11)
        run.font.color.rgb = GREY

    # Three relation labels between the boxes (last one is right-to-left)
    rel_labels = [
        ("produces  →",       0, ACCENT),
        ("represents  →",     1, ACCENT),
        ("←  mastersDataFor", 2, ACCENT),
    ]
    for text, idx, color in rel_labels:
        gap_x = box_xs[idx] + box_w
        lbl = s.shapes.add_textbox(gap_x, y + Inches(0.5), gap, Inches(0.4))
        _set_text(lbl.text_frame, text, size=10, bold=True, color=color, align=PP_ALIGN.CENTER)

    # Sample agent query callout
    callout = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.8),
        prs.slide_width - Inches(1.0), Inches(2.6),
    )
    callout.line.color.rgb = ACCENT
    callout.fill.solid()
    callout.fill.fore_color.rgb = WHITE
    tf = callout.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.2)
    tf.margin_right = Inches(0.3)

    _set_text(tf, "Sample agent query", size=14, bold=True, color=ACCENT)

    p = tf.add_paragraph()
    run = p.add_run()
    run.text = ('\n"Which Activities produce documents whose data is mastered by '
                'SAP S/4 SD?"')
    run.font.size = Pt(15)
    run.font.italic = True
    run.font.color.rgb = NAVY

    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "\nResolved by traversal:"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = NAVY

    for step in [
        "1.  match  ApplicationComponent  where  name = 'SAP S/4 SD'",
        "2.  ← mastersDataFor ─  DataObject  (Customer, SalesOrder, …)",
        "3.  ← represents ─      BusinessDocument  (Sales Order, Invoice, …)",
        "4.  ← produces ─        Activity  (Create Sales Order, Process Invoice, …)",
    ]:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = step
        run.font.size = Pt(11)
        run.font.color.rgb = GREY
        run.font.name = "Consolas"

    p = tf.add_paragraph()
    run = p.add_run()
    run.text = ("\nReturns: 12 Activities, all in the O2C process — answered in <100ms "
                "from the typed graph. The same question against unstructured docs "
                "needs RAG over thousands of pages and still misses the configuration link.")
    run.font.size = Pt(11)
    run.font.color.rgb = GREY

    foot = s.shapes.add_textbox(
        Inches(0.5), Inches(6.6), prs.slide_width - Inches(1.0), Inches(0.4)
    )
    _set_text(
        foot.text_frame,
        "BusinessDocument is the bridge: what users see (a sales order doc) ↔ what the system stores (the SalesOrder data object).",
        size=11, color=ACCENT, bold=True,
    )
    return s


def slide_scenario_provenance(prs):
    s = _blank_slide(prs)
    _add_title_band(s, prs, "Scenario & Provenance", subtitle="Two cross-cutting facts on every instance")
    body = s.shapes.add_textbox(
        Inches(0.5), Inches(1.6), prs.slide_width - Inches(1.0), Inches(5.0)
    )
    _add_bullets(
        body.text_frame,
        [
            "Scenario — every instance lives in at least one named state (draft / active / locked / superseded).",
            "  · scenarioType ∈ {as-is, to-be, variant} — supports parallel target-state design.",
            "  · The runtime supports fork / mutate / promote / diff on Scenarios.",
            "Provenance — every claim records: source system, extracted-by, extracted-at, confidence.",
            "  · Audit-grade lineage: \"this Configuration came from Solman on 2026-04-15 with 0.92 confidence\".",
            "  · Required by the base SHACL shape — missing Provenance fails ingestion outright.",
            "Together: every fact in the graph is dated, sourced, and state-scoped. No more orphan data.",
        ],
        size=15,
    )
    return s


def slide_standards(prs):
    s = _blank_slide(prs)
    _add_title_band(s, prs, "Anchored on industry standards", subtitle="Not invented from scratch")
    rows = [
        ("Architecture", "ArchiMate Application + Technology Layer · LeanIX"),
        ("Business Process", "BPMN 2.0 · APQC Process Classification Framework · Signavio 5-level"),
        ("Implementation / Change", "ITIL 4 · SAP Cloud ALM · SAP TMS"),
        ("Organization", "ArchiMate Business Layer · RACI · TOGAF BIZBOK"),
        ("Scenario / Versioning", "Bitemporal modeling · Git-style branching semantics"),
    ]
    table_shape = s.shapes.add_table(
        len(rows) + 1, 2, Inches(0.5), Inches(1.6),
        prs.slide_width - Inches(1.0), Inches(4.5),
    )
    tbl = table_shape.table
    tbl.columns[0].width = Inches(2.7)
    tbl.columns[1].width = prs.slide_width - Inches(3.7)
    for i, h in enumerate(["Domain", "Anchored on"]):
        cell = tbl.cell(0, i)
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        _set_text(cell.text_frame, h, size=14, bold=True, color=WHITE)
    for r, (dom, anchors) in enumerate(rows, start=1):
        c0 = tbl.cell(r, 0); c1 = tbl.cell(r, 1)
        _set_text(c0.text_frame, dom, size=12, bold=True, color=NAVY)
        _set_text(c1.text_frame, anchors, size=12, color=GREY)
        if r % 2 == 0:
            c0.fill.solid(); c0.fill.fore_color.rgb = LIGHT
            c1.fill.solid(); c1.fill.fore_color.rgb = LIGHT
    return s


def slide_client_usage(prs):
    s = _blank_slide(prs)
    _add_title_band(s, prs, "How clients use it", subtitle="Sub-class, don't fork")
    body = s.shapes.add_textbox(
        Inches(0.5), Inches(1.6), prs.slide_width - Inches(1.0), Inches(5.0)
    )
    _add_bullets(
        body.text_frame,
        [
            "Pin an upper-model version in your project metadata (e.g., v0.2.0).",
            "Sub-class upper classes for client-specific concepts:",
            "    client:CoBuyerProfile  rdfs:subClassOf  sap:Role",
            "Instances carry a client URI scheme:",
            "    https://thbeier.github.io/sap-ontology/client/{clientId}/{domain}/{class}/{id}",
            "Lower-model classes stay private. They flow back upstream only via the contribution process.",
            "Every instance is SHACL-validated at ingestion — Provenance + Scenario are non-negotiable.",
        ],
        size=15,
    )
    return s


def slide_runtime(prs):
    s = _blank_slide(prs)
    _add_title_band(s, prs, "The federation runtime", subtitle="Where the schema meets real client data")
    body = s.shapes.add_textbox(
        Inches(0.5), Inches(1.6), prs.slide_width - Inches(1.0), Inches(2.8)
    )
    _add_bullets(
        body.text_frame,
        [
            "Pipeline:  Extractor  →  Mapper  →  Validator (SHACL)  →  Loader  →  Neo4j",
            "Authoring formats: Excel workbook, CSV directory, raw JSON-LD.",
            "Per-tenant Neo4j databases — no cross-client data leakage.",
            "Scenario operations: fork (clone), mutate (edit), promote (lock + derive impacts), diff.",
            "CLI: sap-ontology-runtime  tenant | pipeline | scenario  ...",
        ],
        size=15,
    )
    pipe = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.1),
        prs.slide_width - Inches(1.0), Inches(1.5),
    )
    pipe.line.color.rgb = ACCENT
    pipe.fill.solid(); pipe.fill.fore_color.rgb = LIGHT
    _set_text(
        pipe.text_frame,
        "  Excel / CSV  →  Mapper  →  SHACL  →  Loader  →  Neo4j (tenant DB)",
        size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
    )
    p = pipe.text_frame.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "\nProprietary repository · pins this schema as a versioned dependency"
    run.font.size = Pt(12)
    run.font.color.rgb = GREY
    return s


def slide_benefits(prs):
    s = _blank_slide(prs)
    _add_title_band(s, prs, "Why clients care", subtitle="Five concrete benefits")
    rows = [
        ("Agent-ready knowledge graph",
         "AI agents can traverse change → config → activity → role without prompt-stuffing the entire SAP wiki."),
        ("Audit-grade lineage",
         "Every fact carries Provenance — source system, timestamp, confidence. No more 'where did this come from?'"),
        ("Scenario branching",
         "Run as-is and to-be in parallel. Fork a Scenario for variant analysis without polluting production state."),
        ("Vendor-neutral by design",
         "Anchored on ArchiMate / BPMN / ITIL — not on a specific tool. Survives Signavio→ARIS or LeanIX→Ardoq migrations."),
        ("Reuse across engagements",
         "The upper model is CC-BY-SA. What you sub-class for one client's pricing flow can inform the next."),
    ]
    table_shape = s.shapes.add_table(
        len(rows), 2, Inches(0.5), Inches(1.5),
        prs.slide_width - Inches(1.0), Inches(5.0),
    )
    tbl = table_shape.table
    tbl.columns[0].width = Inches(2.7)
    tbl.columns[1].width = prs.slide_width - Inches(3.7)
    for r, (label, detail) in enumerate(rows):
        c0 = tbl.cell(r, 0); c1 = tbl.cell(r, 1)
        _set_text(c0.text_frame, label, size=14, bold=True, color=NAVY)
        _set_text(c1.text_frame, detail, size=12, color=GREY)
        if r % 2 == 1:
            c0.fill.solid(); c0.fill.fore_color.rgb = LIGHT
            c1.fill.solid(); c1.fill.fore_color.rgb = LIGHT
    return s


def slide_engagement_model(prs):
    s = _blank_slide(prs)
    _add_title_band(s, prs, "Getting started", subtitle="Engagement model")
    body = s.shapes.add_textbox(
        Inches(0.5), Inches(1.6), prs.slide_width - Inches(1.0), Inches(5.0)
    )
    _add_bullets(
        body.text_frame,
        [
            "1. Discovery (1–2 weeks) — pick a process slice (e.g., O2C). Map source systems to the upper model.",
            "2. Pilot (4–6 weeks) — author a Scenario in Excel, run the runtime end-to-end, validate SHACL, query in Neo4j.",
            "3. Scale — onboard additional Scenarios; sub-class for client-specific concepts; integrate AI agents.",
            "Governance: Deloitte SAP Operate + Methods & Tools steer the upper model. External PRs welcome.",
            "License: schema CC-BY-SA 4.0. Mapper configurations and runtime are Deloitte proprietary.",
        ],
        size=15,
    )
    return s


def slide_next_steps(prs):
    s = _blank_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY

    title = s.shapes.add_textbox(
        Inches(0.8), Inches(2.0), prs.slide_width - Inches(1.6), Inches(1.0)
    )
    _set_text(title.text_frame, "Let's pilot it.", size=44, bold=True, color=WHITE)

    sub = s.shapes.add_textbox(
        Inches(0.8), Inches(3.4), prs.slide_width - Inches(1.6), Inches(2.5)
    )
    _add_bullets(
        sub.text_frame,
        [
            "Pick a process — O2C, P2P, R2R, or your highest-pain area.",
            "We'll spin up a Scenario in 4–6 weeks and run the first agent queries against your data.",
            "Discussion: which process? which source systems? which target Neo4j environment?",
        ],
        size=18,
    )
    for p in sub.text_frame.paragraphs:
        for r in p.runs:
            r.font.color.rgb = WHITE
    return s


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    builders = [
        slide_title,
        slide_problem,
        slide_what_it_is,
        slide_five_domains,
        slide_class_catalog,
        slide_inter_domain,
        slide_artifact_bridge,
        slide_scenario_provenance,
        slide_standards,
        slide_client_usage,
        slide_runtime,
        slide_benefits,
        slide_engagement_model,
        slide_next_steps,
    ]
    total = len(builders)
    for i, build in enumerate(builders, start=1):
        slide = build(prs)
        if i not in (1, total):  # skip footer on title + closing
            _add_footer(slide, prs, i, total)

    prs.save(DECK_PATH)
    print(f"Wrote {DECK_PATH} ({total} slides)")


if __name__ == "__main__":
    main()
